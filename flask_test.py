# app.py

from flask import Flask, request, redirect, render_template
from pptx import Presentation
from pptx.util import Inches
from os import listdir
import os

app = Flask(__name__)

@app.route('/')
def index():
    return render_template('flask_test.html')

UPLOAD_FOLDER = 'static/uploads'
app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER

@app.route("/upload", methods=["POST"])
def upload():
    file = request.files["file"]
    if file:
        file.save(os.path.join(app.config['UPLOAD_FOLDER'], file.filename))
        return redirect("/")
    else:
        return "File not uploaded"


@app.route('/convert', methods=['POST'])
def convert():
    directory = request.form['directory']
    filenames = [f for f in listdir(directory) if f.endswith('.jpg')]

    prs = Presentation()
    for filename in filenames:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        left = top = Inches(1)
        height = Inches(5)
        pic = slide.shapes.add_picture(directory + '/' + filename, left, top, height=height)

    prs.save('converted.pptx')
    return 'PPTX file has been converted.'

from flask import Flask, send_file

app = Flask(__name__)

@app.route("/download_pptx")
def download_pptx():
    filename = "presentation.pptx"
    return send_file(filename, as_attachment=True)


if __name__ == '__main__':
    app.run(debug=True)
